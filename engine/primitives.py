"""
Drawing Primitives — مذكرتي Pro v17.1
Low-level, deterministic shape/text builders.

FIXED BUGS vs v17.0:
- gradient_fill() now inserts gradFill in CORRECT OOXML position (before <a:ln>)
- shadow() now inserts effectLst in CORRECT position (after <a:ln>)
- _sort_spPr() enforces strict OOXML child order on every shape
- set_solid_alpha() exposed (no duplication with slides.py)
"""
from __future__ import annotations

from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# Slide dimensions (cm) — 16:9
W, H = 33.867, 19.05

# OOXML spPr required child order
_SPPR_ORDER = [
    qn('a:xfrm'), qn('a:prstGeom'), qn('a:custGeom'),
    qn('a:noFill'), qn('a:solidFill'), qn('a:gradFill'),
    qn('a:blipFill'), qn('a:pattFill'), qn('a:grpFill'),
    qn('a:ln'),
    qn('a:effectLst'), qn('a:effectDag'),
    qn('a:scene3d'), qn('a:sp3d'), qn('a:extLst'),
]
_SPPR_RANK = {tag: i for i, tag in enumerate(_SPPR_ORDER)}


def _sort_spPr(spPr) -> None:
    """Reorder <p:spPr> children to comply with OOXML schema."""
    children = list(spPr)
    children.sort(key=lambda el: _SPPR_RANK.get(el.tag, 99))
    for child in children:
        spPr.remove(child)
    for child in children:
        spPr.append(child)


def _get_spPr(shape):
    return shape._element.find(qn('p:spPr'))


# ── Unit helpers ─────────────────────────────────────────────────────
def cm(v: float) -> int:
    return int(Cm(v))

def pt(v: float) -> int:
    return int(Pt(v))


# ── Shape builders ───────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill: RGBColor, line=None, line_w=0.5):
    if w <= 0 or h <= 0:
        return None
    s = slide.shapes.add_shape(1, cm(x), cm(y), cm(w), cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = pt(line_w)
    else:
        s.line.fill.background()
    return s


def rrect(slide, x, y, w, h, fill: RGBColor, radius_pct=8, line=None, line_w=0.5):
    if w <= 0 or h <= 0:
        return None
    s = slide.shapes.add_shape(5, cm(x), cm(y), cm(w), cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = pt(line_w)
    else:
        s.line.fill.background()
    try:
        adj = s.adjustments
        if adj and len(adj) > 0:
            adj[0] = max(0, min(50, radius_pct)) * 1000
    except Exception:
        pass
    return s


def oval(slide, x, y, w, h, fill: RGBColor, alpha=100):
    if w <= 0 or h <= 0:
        return None
    s = slide.shapes.add_shape(9, cm(x), cm(y), cm(w), cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if alpha < 100:
        set_solid_alpha(s, alpha)
    return s


def bg(slide, color: RGBColor):
    rect(slide, 0, 0, W, H, color)

def hline(slide, x, y, w, color: RGBColor, thickness=0.08):
    rect(slide, x, y, w, thickness, color)

def vline(slide, x, y, h2, color: RGBColor, thickness=0.08):
    rect(slide, x, y, thickness, h2, color)


# ── XML fill helpers ─────────────────────────────────────────────────
def set_solid_alpha(shape, alpha_pct: int) -> None:
    """Set transparency on a solidFill shape (0=transparent, 100=opaque)."""
    try:
        spPr = _get_spPr(shape)
        srgb = spPr.find('.//' + qn('a:srgbClr'))
        if srgb is not None:
            for e in srgb.findall(qn('a:alpha')):
                srgb.remove(e)
            alp = etree.SubElement(srgb, qn('a:alpha'))
            alp.set('val', str(int(alpha_pct * 1000)))
    except Exception:
        pass


def gradient_fill(shape, c1: str, c2: str, angle: float = 90) -> None:
    """
    Apply linear gradient via XML, with correct OOXML element ordering.
    BUG FIX: previously appended gradFill AFTER <a:ln> — now sorted correctly.
    """
    try:
        spPr = _get_spPr(shape)
        # Remove all existing fill variants
        for tag in [qn('a:solidFill'), qn('a:gradFill'), qn('a:noFill'),
                    qn('a:pattFill'), qn('a:blipFill'), qn('a:grpFill')]:
            for el in spPr.findall(tag):
                spPr.remove(el)

        # Build gradFill
        grad = etree.Element(qn('a:gradFill'))
        gsLst = etree.SubElement(grad, qn('a:gsLst'))

        gs0 = etree.SubElement(gsLst, qn('a:gs'))
        gs0.set('pos', '0')
        etree.SubElement(gs0, qn('a:srgbClr')).set('val', c1.lstrip('#'))

        gs1 = etree.SubElement(gsLst, qn('a:gs'))
        gs1.set('pos', '100000')
        etree.SubElement(gs1, qn('a:srgbClr')).set('val', c2.lstrip('#'))

        lin = etree.SubElement(grad, qn('a:lin'))
        lin.set('ang', str(int(angle * 60000)))
        lin.set('scaled', '0')

        spPr.append(grad)
        _sort_spPr(spPr)  # ← enforce correct order
    except Exception:
        pass


def gradient_rect(slide, x, y, w, h, c1: str, c2: str, angle=0):
    c1h = c1.lstrip('#')
    fill_color = RGBColor(int(c1h[0:2], 16), int(c1h[2:4], 16), int(c1h[4:6], 16))
    s = rect(slide, x, y, w, h, fill_color)
    if s:
        gradient_fill(s, c1, c2, angle)
    return s


def shadow(shape, blur=16, dist=5, angle=135, alpha=0.22, color="000000") -> None:
    """
    Add outer drop shadow via XML.
    BUG FIX: effectLst now inserted in correct OOXML position (after <a:ln>).
    """
    try:
        spPr = _get_spPr(shape)
        for old in spPr.findall(qn('a:effectLst')):
            spPr.remove(old)

        eLst = etree.Element(qn('a:effectLst'))
        shdw = etree.SubElement(eLst, qn('a:outerShdw'))
        shdw.set('blurRad', str(int(blur * 12700)))
        shdw.set('dist', str(int(dist * 12700)))
        shdw.set('dir', str(int(angle * 60000)))
        shdw.set('algn', 'tl')
        srgb = etree.SubElement(shdw, qn('a:srgbClr'))
        srgb.set('val', color.lstrip('#'))
        alp = etree.SubElement(srgb, qn('a:alpha'))
        alp.set('val', str(int(alpha * 100000)))

        spPr.append(eLst)
        _sort_spPr(spPr)  # ← enforce correct order
    except Exception:
        pass


# ── Text ─────────────────────────────────────────────────────────────
def txt(slide, text: str, x, y, w, h,
        font="Cairo", size=14, bold=False, italic=False,
        color: RGBColor | None = None,
        align=PP_ALIGN.RIGHT,
        margin=0.1, rtl=True, spacing=None,
        vcenter=True, line_spacing=1.15):
    """
    نص احترافي باستخدام Shape مع توسيط عمودي حقيقي.
    vcenter=True → MSO_ANCHOR.MIDDLE (يعمل في PowerPoint وLibreOffice)
    line_spacing → ارتفاع السطر النسبي
    """
    if not text or w <= 0 or h <= 0:
        return None
    try:
        from pptx.enum.text import MSO_ANCHOR
        sh = slide.shapes.add_shape(1, cm(x), cm(y), cm(w), cm(h))
        sh.fill.background()
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True
        tf.margin_left   = cm(margin)
        tf.margin_right  = cm(margin)
        tf.margin_top    = cm(0.04)
        tf.margin_bottom = cm(0.04)

        if vcenter:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        else:
            tf.vertical_anchor = MSO_ANCHOR.TOP

        p = tf.paragraphs[0]
        p.alignment = align

        # ارتفاع السطر
        try:
            from pptx.oxml.ns import qn as _qn
            pPr = p._p.get_or_add_pPr()
            # أزل lnSpc القديم إن وُجد
            for old in pPr.findall(_qn('a:lnSpc')):
                pPr.remove(old)
            lnSpc = etree.SubElement(pPr, _qn('a:lnSpc'))
            spcPct = etree.SubElement(lnSpc, _qn('a:spcPct'))
            spcPct.set('val', str(int(line_spacing * 100000)))
        except Exception:
            pass

        run = p.add_run()
        run.text = str(text)
        run.font.name   = font
        run.font.size   = Pt(size)
        run.font.bold   = bold
        run.font.italic = italic
        if color:
            run.font.color.rgb = color
        return sh
    except Exception:
        # fallback إلى textbox
        tb = slide.shapes.add_textbox(cm(x), cm(y), cm(w), cm(h))
        tb.word_wrap = True
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.name   = font
        run.font.size   = Pt(size)
        run.font.bold   = bold
        run.font.italic = italic
        if color:
            run.font.color.rgb = color
        return tb


def txt2(slide, label: str, value: str, x, y, w, h,
         font="Cairo", label_size=10, value_size=13,
         label_color: RGBColor | None = None,
         value_color: RGBColor | None = None,
         align=PP_ALIGN.RIGHT, margin=0.1):
    """
    نص بسطرين: تسمية + قيمة مع توسيط عمودي.
    مثالي لبطاقات المعلومات.
    """
    if w <= 0 or h <= 0: return None
    try:
        from pptx.enum.text import MSO_ANCHOR
        sh = slide.shapes.add_shape(1, cm(x), cm(y), cm(w), cm(h))
        sh.fill.background()
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True
        tf.margin_left   = cm(margin)
        tf.margin_right  = cm(margin)
        tf.margin_top    = cm(0.04)
        tf.margin_bottom = cm(0.04)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p1 = tf.paragraphs[0]
        p1.alignment = align
        r1 = p1.add_run()
        r1.text = str(label)
        r1.font.name  = font
        r1.font.size  = Pt(label_size)
        r1.font.bold  = True
        if label_color: r1.font.color.rgb = label_color

        p2 = tf.add_paragraph()
        p2.alignment = align
        r2 = p2.add_run()
        r2.text = str(value)
        r2.font.name  = font
        r2.font.size  = Pt(value_size)
        r2.font.bold  = False
        if value_color: r2.font.color.rgb = value_color
        return sh
    except Exception:
        return None


def blank_slide(prs):
    """Add a completely blank slide (layout 6 = Blank)."""
    return prs.slides.add_slide(prs.slide_layouts[6])


# ── Advanced Visual Primitives ────────────────────────────────────────

def glow(shape, color: str = "C6A03C", radius: int = 20, alpha: float = 0.4) -> None:
    """Add a glow effect (outerShdw with zero distance = glow)."""
    try:
        spPr = _get_spPr(shape)
        for old in spPr.findall(qn('a:effectLst')):
            spPr.remove(old)
        eLst = etree.Element(qn('a:effectLst'))
        g = etree.SubElement(eLst, qn('a:outerShdw'))
        g.set('blurRad', str(int(radius * 12700)))
        g.set('dist', '0')
        g.set('dir', '0')
        g.set('algn', 'ctr')
        srgb = etree.SubElement(g, qn('a:srgbClr'))
        srgb.set('val', color.lstrip('#'))
        alp = etree.SubElement(srgb, qn('a:alpha'))
        alp.set('val', str(int(alpha * 100000)))
        spPr.append(eLst)
        _sort_spPr(spPr)
    except Exception:
        pass


def multi_stop_gradient(shape, stops: list[tuple[int, str]], angle: float = 90) -> None:
    """
    Apply a multi-stop gradient.
    stops = [(pos_pct, '#RRGGBB'), ...]  e.g. [(0,'#07172F'),(50,'#1A3A6E'),(100,'#07172F')]
    """
    try:
        spPr = _get_spPr(shape)
        for tag in [qn('a:solidFill'), qn('a:gradFill'), qn('a:noFill'),
                    qn('a:pattFill'), qn('a:blipFill'), qn('a:grpFill')]:
            for el in spPr.findall(tag):
                spPr.remove(el)
        grad = etree.Element(qn('a:gradFill'))
        gsLst = etree.SubElement(grad, qn('a:gsLst'))
        for pos_pct, hex_color in stops:
            gs = etree.SubElement(gsLst, qn('a:gs'))
            gs.set('pos', str(int(pos_pct * 1000)))
            etree.SubElement(gs, qn('a:srgbClr')).set('val', hex_color.lstrip('#'))
        lin = etree.SubElement(grad, qn('a:lin'))
        lin.set('ang', str(int(angle * 60000)))
        lin.set('scaled', '0')
        spPr.append(grad)
        _sort_spPr(spPr)
    except Exception:
        pass


def gradient_oval(slide, x, y, w, h, c1: str, c2: str, angle=90, alpha=100):
    """Oval with gradient fill."""
    if w <= 0 or h <= 0:
        return None
    c1h = c1.lstrip('#')
    fill_color = RGBColor(int(c1h[0:2], 16), int(c1h[2:4], 16), int(c1h[4:6], 16))
    s = slide.shapes.add_shape(9, cm(x), cm(y), cm(w), cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.fill.background()
    gradient_fill(s, c1, c2, angle)
    if alpha < 100:
        set_solid_alpha(s, alpha)
    return s


def triangle(slide, x, y, w, h, fill: RGBColor, pointing='up'):
    """Equilateral-ish triangle shape."""
    if w <= 0 or h <= 0:
        return None
    # Use right-triangle preset (shape 6) rotated for pointing direction
    shape_id = 6  # rtTriangle
    s = slide.shapes.add_shape(shape_id, cm(x), cm(y), cm(w), cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if pointing == 'down':
        s.rotation = 180
    elif pointing == 'left':
        s.rotation = 90
    elif pointing == 'right':
        s.rotation = 270
    return s


def diamond(slide, x, y, w, h, fill: RGBColor, alpha=100):
    """Diamond shape."""
    if w <= 0 or h <= 0:
        return None
    s = slide.shapes.add_shape(4, cm(x), cm(y), cm(w), cm(h))  # shape 4 = diamond
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if alpha < 100:
        set_solid_alpha(s, alpha)
    return s


def hexagon(slide, x, y, w, h, fill: RGBColor, alpha=100):
    """Hexagon shape."""
    if w <= 0 or h <= 0:
        return None
    s = slide.shapes.add_shape(10, cm(x), cm(y), cm(w), cm(h))  # shape 10 = hexagon
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if alpha < 100:
        set_solid_alpha(s, alpha)
    return s


def arc_progress(slide, x, y, size, fill: RGBColor, bg_color: RGBColor,
                 thickness=0.4) -> None:
    """Simulated progress ring using two arcs (outer ring + inner mask)."""
    # Outer ring background
    outer = oval(slide, x, y, size, size, bg_color, alpha=30)
    # Inner circle mask (creates ring effect)
    inner_offset = thickness
    inner_s = size - 2 * inner_offset
    oval(slide, x + inner_offset, y + inner_offset, inner_s, inner_s, bg_color, alpha=80)


def decorative_dots(slide, x, y, cols, rows, dot_size, gap, color: RGBColor, alpha=20):
    """Grid of decorative dots."""
    for r in range(rows):
        for c in range(cols):
            dx = x + c * (dot_size + gap)
            dy = y + r * (dot_size + gap)
            o = oval(slide, dx, dy, dot_size, dot_size, color)
            if o and alpha < 100:
                set_solid_alpha(o, alpha)


def wave_rect(slide, x, y, w, h, fill: RGBColor, wavy_top=True):
    """Rectangle with rounded top (simulates wave). Uses rrect with high radius."""
    if wavy_top:
        return rrect(slide, x, y, w, h, fill, radius_pct=12)
    return rect(slide, x, y, w, h, fill)


def badge(slide, x, y, w, h, fill_c1: str, fill_c2: str, label: str,
          font="Cairo", font_size=11, text_color: RGBColor = None, T=None):
    """Pill-shaped badge with gradient and centered label."""
    b = rrect(slide, x, y, w, h, RGBColor(0xC6, 0xA0, 0x3C), radius_pct=50)
    if b:
        gradient_fill(b, fill_c1, fill_c2, angle=0)
    if text_color is None and T is not None:
        text_color = T.text_dark_rgb
    txt(slide, label, x, y, w, h,
        font=font, size=font_size, bold=True,
        color=text_color, align=PP_ALIGN.CENTER, rtl=True)
    return b


def icon_circle(slide, x, y, size, bg_c1: str, bg_c2: str,
                icon: str, icon_size=20, T=None):
    """Circle with gradient bg + centered emoji/icon."""
    c = oval(slide, x, y, size, size,
             RGBColor(int(bg_c1.lstrip('#')[0:2], 16),
                      int(bg_c1.lstrip('#')[2:4], 16),
                      int(bg_c1.lstrip('#')[4:6], 16)))
    if c:
        gradient_fill(c, bg_c1, bg_c2, angle=135)
    txt(slide, icon, x, y, size, size,
        font="Calibri", size=icon_size, bold=False,
        color=T.text_dark_rgb if T else RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER, rtl=False)
    return c


def number_badge(slide, x, y, size, num: int | str, T):
    """Circular number badge with accent gradient."""
    c = oval(slide, x, y, size, size, T.accent_rgb)
    if c:
        gradient_fill(c, T.accent_grad1, T.accent_grad2, 135)
        shadow(c, blur=10, dist=3, alpha=0.35)
    txt(slide, str(num), x, y, size, size,
        font="Calibri", size=max(8, int(size * 18)), bold=True,
        color=T.text_dark_rgb, align=PP_ALIGN.CENTER, rtl=False)
    return c


def divider(slide, x, y, w, T, style='gradient'):
    """Decorative divider line."""
    if style == 'gradient':
        d = rect(slide, x, y, w, 0.06, T.accent_rgb)
        if d:
            multi_stop_gradient(d, [(0, T.bg), (50, T.accent), (100, T.bg)], angle=0)
    elif style == 'double':
        rect(slide, x, y, w, 0.05, T.accent_rgb)
        rect(slide, x + w * 0.05, y + 0.12, w * 0.9, 0.03, T.muted_rgb)
    else:
        rect(slide, x, y, w, 0.06, T.accent_rgb)


def card_3d(slide, x, y, w, h, T, radius=10):
    """Card with shadow + subtle gradient for 3D feel."""
    # Shadow layer (slightly offset)
    shadow_s = rrect(slide, x + 0.15, y + 0.2, w, h, T.bg_rgb, radius_pct=radius)
    if shadow_s:
        set_solid_alpha(shadow_s, 40)

    # Main card
    c = rrect(slide, x, y, w, h, T.card_rgb, radius_pct=radius)
    if c:
        multi_stop_gradient(c, [(0, T.card), (100, T.bg2)], angle=135)
        shadow(c, blur=18, dist=5, alpha=0.4)
    return c


def slide_number(slide, num: int, total: int, T):
    """Slide number indicator bottom-right."""
    label = f"{num} / {total}"
    txt(slide, label, W - 3.5, H - 0.55, 3.2, 0.45,
        font="Calibri", size=9, bold=False,
        color=T.muted_rgb, align=PP_ALIGN.LEFT, rtl=False)


def watermark(slide, text: str, T):
    """Subtle watermark bottom-left."""
    txt(slide, text, 0.4, H - 0.55, 6.0, 0.45,
        font="Calibri", size=8, bold=False,
        color=T.muted_rgb, align=PP_ALIGN.RIGHT, rtl=False)


def section_tag(slide, label: str, x, y, T):
    """Small colored tag/label."""
    w, h = 3.5, 0.52
    b = rrect(slide, x, y, w, h, T.accent_rgb, radius_pct=50)
    if b:
        gradient_fill(b, T.accent_grad1, T.accent_grad2, 0)
    txt(slide, label, x, y, w, h,
        font="Cairo", size=10, bold=True,
        color=T.text_dark_rgb, align=PP_ALIGN.CENTER, rtl=True)


# ══════════════════════════════════════════════════════════════════════
# DESIGN INTELLIGENCE LAYER — v28
# Smart layout utilities that respond to content
# ══════════════════════════════════════════════════════════════════════

def _smart_font_size(text: str, base: float, min_s: float, max_s: float,
                      area_w: float, area_h: float,
                      chars_per_pt: float = 0.065) -> float:
    """
    Compute font size that fits `text` inside (area_w × area_h).
    chars_per_pt = approximate chars per pt width per cm of width.
    """
    if not text:
        return base
    length = len(text)
    # Very short → go large; very long → go small
    if length <= 20:
        factor = 1.18
    elif length <= 40:
        factor = 1.0
    elif length <= 70:
        factor = 0.88
    elif length <= 110:
        factor = 0.76
    else:
        factor = 0.66
    size = base * factor
    return max(min_s, min(max_s, size))


def smart_title(slide, text: str, x, y, w, h, T, font="Cairo",
                base_size=30, min_s=18, max_s=38, rtl=True, vcenter=True):
    """
    Dominant section title with smart sizing and accent underline.
    """
    from pptx.enum.text import PP_ALIGN
    fs = _smart_font_size(text, base_size, min_s, max_s, w, h)
    title_h = h * 0.72 if h > 1.0 else h
    t = txt(slide, text, x, y, w, title_h,
            font=font, size=fs, bold=True,
            color=T.text_light_rgb, align=PP_ALIGN.RIGHT,
            rtl=rtl, vcenter=vcenter, line_spacing=1.05)
    return t


def accent_pill(slide, text: str, x, y, T, font="Cairo", size=10.5):
    """Small accent-colored pill label."""
    from pptx.enum.text import PP_ALIGN
    w = max(2.8, len(text) * 0.18 + 0.6)
    h = 0.44
    b = rrect(slide, x, y, w, h, T.accent_rgb, radius_pct=50)
    if b:
        gradient_fill(b, T.accent_grad1, T.accent_grad2, 0)
    txt(slide, text, x, y, w, h,
        font=font, size=size, bold=True,
        color=T.text_dark_rgb, align=PP_ALIGN.CENTER, rtl=True)
    return w, h


def premium_card(slide, x, y, w, h, T, radius=12, depth=True, glow_on=True):
    """
    Premium card with layered shadow, gradient, optional glow.
    Returns the main card shape.
    """
    if depth:
        # Deep shadow base
        sh = rrect(slide, x + 0.18, y + 0.24, w, h, T.bg_rgb, radius_pct=radius)
        if sh:
            set_solid_alpha(sh, 35)
    c = rrect(slide, x, y, w, h, T.card_rgb, radius_pct=radius)
    if c:
        multi_stop_gradient(c, [(0, T.card), (60, T.bg2), (100, T.bg)], 135)
        shadow(c, blur=22, dist=6, alpha=0.44)
        if glow_on:
            glow(c, T.accent.lstrip('#'), radius=16, alpha=0.07)
    return c


def card_with_accent_top(slide, x, y, w, h, T, radius=12, bar_h=0.38):
    """Card with colored accent bar on top."""
    c = premium_card(slide, x, y, w, h, T, radius=radius)
    bar = rrect(slide, x, y, w, bar_h, T.accent_rgb, radius_pct=0)
    if bar:
        multi_stop_gradient(bar, [(0, T.accent2), (50, T.accent), (100, T.accent2)], 0)
        glow(bar, T.accent.lstrip('#'), radius=10, alpha=0.22)
    return c


def card_with_accent_side(slide, x, y, w, h, T, radius=12, bar_w=0.26):
    """Card with colored accent bar on right side (RTL primary)."""
    c = premium_card(slide, x, y, w, h, T, radius=radius)
    bar = rrect(slide, x + w - bar_w, y, bar_w, h, T.accent_rgb, radius_pct=0)
    if bar:
        gradient_fill(bar, T.accent_grad1, T.accent_grad2, 90)
    return c


def kpi_card(slide, x, y, w, h, T, value: str, label: str,
             unit: str = '', font="Cairo"):
    """
    Premium KPI/stat card: giant centered value + label below.
    Auto-scales value font based on length.
    """
    from pptx.enum.text import PP_ALIGN
    # Card base
    c = rrect(slide, x, y, w, h, T.card_rgb, radius_pct=14)
    if c:
        multi_stop_gradient(c, [(0, T.bg2), (50, T.card), (100, T.bg2)], 135)
        shadow(c, blur=20, dist=6, alpha=0.45)

    # Accent top stripe
    tp = rrect(slide, x, y, w, 0.32, T.accent_rgb, radius_pct=0)
    if tp:
        multi_stop_gradient(tp, [(0, T.accent2), (50, T.accent), (100, T.accent2)], 0)
        glow(tp, T.accent.lstrip('#'), radius=8, alpha=0.28)

    # Bottom pulse bar
    bp = rrect(slide, x, y + h - 0.18, w, 0.18, T.accent_rgb, radius_pct=0)
    if bp:
        set_solid_alpha(bp, 30)

    # Value — giant, centered
    vlen = len(str(value))
    vs = 46 if vlen <= 2 else 38 if vlen <= 4 else 28 if vlen <= 7 else 22
    txt(slide, str(value), x + 0.12, y + 0.32, w - 0.24, h * 0.50,
        font="Calibri", size=vs, bold=True,
        color=T.accent_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)

    # Unit badge (if present)
    if unit:
        ub = rrect(slide, x + w / 2 - 1.6, y + h * 0.53 + 0.06, 3.2, 0.42,
                   T.bg_rgb, radius_pct=40)
        if ub:
            set_solid_alpha(ub, 55)
        txt(slide, unit, x + w / 2 - 1.6, y + h * 0.53 + 0.06, 3.2, 0.42,
            font=font, size=9.5, bold=False,
            color=T.muted_rgb, align=PP_ALIGN.CENTER, rtl=True, vcenter=True)

    # Divider
    hline(slide, x + w * 0.14, y + h * 0.71, w * 0.72, T.muted_rgb, thickness=0.04)

    # Label
    txt(slide, label, x + 0.12, y + h * 0.73, w - 0.24, h * 0.25,
        font=font, size=max(10, min(13, h * 5.5)), bold=False,
        color=T.text_light_rgb, align=PP_ALIGN.CENTER, rtl=True, vcenter=True)


def result_row(slide, x, y, w, h, T, text: str, index: int,
               font="Cairo", highlight=False):
    """
    Premium result list row with smart font sizing.
    index: 1-based number shown in badge on right.
    highlight: makes row visually stronger (for key results).
    """
    from pptx.enum.text import PP_ALIGN
    even = (index % 2 == 0)

    # Row background
    rw = rrect(slide, x, y, w, h,
               T.card_rgb if not even else T.bg2_rgb,
               radius_pct=10)
    if rw:
        stops = [(0, T.card), (100, T.bg2)] if not even else [(0, T.bg2), (100, T.card)]
        multi_stop_gradient(rw, stops, 0)
        if highlight:
            shadow(rw, blur=12, dist=3, alpha=0.32)
            glow(rw, T.accent.lstrip('#'), radius=12, alpha=0.06)
        else:
            shadow(rw, blur=5, dist=2, alpha=0.16)

    # Accent side bar (fades with index for visual rhythm)
    alpha_bar = max(22, 62 - index * 6)
    bar = rect(slide, x + w - 0.28, y, 0.28, h, T.accent_rgb)
    if bar:
        gradient_fill(bar, T.accent_grad1, T.accent_grad2, 90)
        set_solid_alpha(bar, alpha_bar)

    # Number badge
    nd = min(0.68, h * 0.72)
    nb_x = x + w - 1.1 - nd
    nb_y = y + (h - nd) / 2
    nb_c = oval(slide, nb_x, nb_y, nd, nd, T.accent_rgb)
    if nb_c:
        multi_stop_gradient(nb_c, [(0, T.accent), (100, T.accent2)], 135)
        shadow(nb_c, blur=8, dist=2, alpha=0.3)
    txt(slide, str(index), nb_x, nb_y, nd, nd,
        font="Calibri", size=max(9, int(nd * 11)), bold=True,
        color=T.text_dark_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)

    # Content text — smart sizing
    text_w = w - nd - 1.5
    fs = _smart_font_size(text, 13.5, 11, 15.5, text_w, h)
    txt(slide, text, x + 0.25, y, text_w, h,
        font=font, size=fs, bold=highlight,
        color=T.text_light_rgb if not highlight else T.accent_rgb,
        align=PP_ALIGN.RIGHT, rtl=True, vcenter=True, line_spacing=1.2)


def premium_header(slide, T, title: str, subtitle: str = '',
                   slide_num: int = None, total: int = 13,
                   accent_side='right', font="Cairo"):
    """
    Premium header with:
    - Dominant gradient background
    - Strong title with smart sizing
    - Subtle subtitle
    - Slide counter badge
    - Multi-layer accent lines
    """
    from pptx.enum.text import PP_ALIGN
    HDR_H = 3.0

    # Main header background — deep gradient
    gradient_rect(slide, 0, 0, W, HDR_H, T.grad2, T.grad1, angle=130)

    # Layered accent lines at bottom
    al1 = rect(slide, 0, HDR_H - 0.26, W, 0.26, T.accent_rgb)
    if al1:
        multi_stop_gradient(al1, [(0, T.bg), (35, T.accent2), (50, T.accent),
                                   (65, T.accent2), (100, T.bg)], 0)
    rect(slide, 0, HDR_H - 0.32, W, 0.06, T.muted_rgb)
    rect(slide, 0, HDR_H - 0.06, W, 0.06, T.bg_rgb)

    # Accent vertical bar
    if accent_side == 'right':
        av = rect(slide, W - 0.56, 0, 0.56, HDR_H, T.accent_rgb)
    else:
        av = rect(slide, 0, 0, 0.56, HDR_H, T.accent_rgb)
    if av:
        gradient_fill(av, T.accent_grad1, T.accent_grad2, 90)

    # Decorative background circle
    oval(slide, W - 5.5, -2.5, 8, 8, T.accent_rgb, alpha=8)

    # Slide number badge
    if slide_num is not None:
        nb_s = 0.78
        nb_x = 1.05
        nb_y = (HDR_H - nb_s) / 2
        nb_c = oval(slide, nb_x, nb_y, nb_s, nb_s, T.accent_rgb)
        if nb_c:
            multi_stop_gradient(nb_c, [(0, T.accent_grad1), (100, T.accent_grad2)], 135)
            shadow(nb_c, blur=10, dist=3, alpha=0.38)
        txt(slide, str(slide_num), nb_x, nb_y, nb_s, nb_s,
            font="Calibri", size=15, bold=True,
            color=T.text_dark_rgb, align=PP_ALIGN.CENTER, rtl=False, vcenter=True)
        txt(slide, f"/{total}", nb_x + nb_s, nb_y + nb_s * 0.32, 0.85, nb_s * 0.38,
            font="Calibri", size=8, bold=False,
            color=T.muted_rgb, align=PP_ALIGN.LEFT, rtl=False, vcenter=True)
        title_x = nb_x + nb_s + 0.9
    else:
        title_x = 0.72

    title_w = W - title_x - 0.72
    fs_title = _smart_font_size(title, 30, 20, 34, title_w, HDR_H * 0.65)
    txt(slide, title, title_x, 0.18, title_w, HDR_H * 0.63,
        font=font, size=fs_title, bold=True,
        color=T.text_light_rgb, align=PP_ALIGN.RIGHT,
        rtl=True, vcenter=True, line_spacing=1.05)

    if subtitle:
        fs_sub = min(14.5, max(11, fs_title * 0.44))
        txt(slide, subtitle, title_x, HDR_H * 0.63, title_w, HDR_H * 0.33,
            font=font, size=fs_sub, bold=False, italic=True,
            color=T.muted_rgb, align=PP_ALIGN.RIGHT,
            rtl=True, vcenter=True, line_spacing=1.0)

    return HDR_H


def section_divider_line(slide, x, y, w, T):
    """Triple-layer decorative divider."""
    d1 = rect(slide, x, y, w, 0.07, T.accent_rgb)
    if d1:
        multi_stop_gradient(d1, [(0, T.bg2), (50, T.accent), (100, T.bg2)], 0)
    rect(slide, x + w * 0.08, y + 0.1, w * 0.84, 0.03, T.muted_rgb)


def two_col_layout(n_items):
    """Return (cols, rows) for n items, preferring 2-col layout when n>3."""
    if n_items <= 3:
        return n_items, 1
    elif n_items <= 6:
        return 2, (n_items + 1) // 2
    else:
        return 3, (n_items + 2) // 3


def adaptive_body_size(text: str, container_h: float,
                        base=13.5, min_s=10.5, max_s=16.0) -> float:
    """Scale body text to fill a container height comfortably."""
    n_words = len(text.split())
    if n_words <= 10:
        factor = 1.15
    elif n_words <= 20:
        factor = 1.0
    elif n_words <= 35:
        factor = 0.88
    elif n_words <= 55:
        factor = 0.76
    else:
        factor = 0.65
    size = base * factor
    # Also constrain to height
    h_factor = container_h * 4.5
    size = min(size, h_factor)
    return max(min_s, min(max_s, size))


def premium_bg(slide, T, style='a'):
    """
    Enhanced background with depth layers and ambient shapes.
    Styles: 'a' (radial), 'b' (diagonal), 'c' (corner), 'd' (concentric)
    """
    bg(slide, T.bg_rgb)
    angle_map = {'a': 135, 'b': 160, 'c': 90, 'd': 45}
    gradient_rect(slide, 0, 0, W, H, T.grad1, T.grad2,
                  angle=angle_map.get(style, 135))

    if style == 'a':
        oval(slide, -4, -4, 13, 13, T.accent_rgb, alpha=5)
        oval(slide, W - 10, H - 9, 15, 15, T.bg2_rgb, alpha=42)
        oval(slide, W - 7, -1, 9, 9, T.accent_rgb, alpha=4)
        decorative_dots(slide, 1.2, H - 4.5, 5, 3, 0.16, 0.42, T.accent_rgb, alpha=11)
    elif style == 'b':
        diamond(slide, W - 7.5, -2.5, 6.5, 6.5, T.accent_rgb, alpha=6)
        diamond(slide, -1.5, H - 5, 5, 5, T.accent_rgb, alpha=5)
        hexagon(slide, W - 5, H * 0.3, 3.0, 3.0, T.accent_rgb, alpha=7)
        decorative_dots(slide, 1.0, 1.8, 4, 4, 0.15, 0.36, T.accent_rgb, alpha=9)
        oval(slide, W * 0.35, -3, 8, 8, T.accent_rgb, alpha=3)
    elif style == 'c':
        oval(slide, -5, -4, 13, 13, T.accent_rgb, alpha=4)
        oval(slide, W - 11, H - 10, 16, 16, T.accent_rgb, alpha=4)
        oval(slide, W - 7, -3, 10, 10, T.bg2_rgb, alpha=38)
        decorative_dots(slide, W - 7, 1.5, 4, 5, 0.14, 0.35, T.accent_rgb, alpha=10)
        oval(slide, -2, H * 0.4, 6, 6, T.bg2_rgb, alpha=22)
    elif style == 'd':
        for r, a in [(28, 3), (22, 4), (16, 5), (10, 7), (6, 9)]:
            oval(slide, W / 2 - r / 2, H / 2 - r / 2, r, r, T.accent_rgb, alpha=a)
        decorative_dots(slide, 1.8, H - 4.2, 5, 2, 0.18, 0.44, T.accent_rgb, alpha=11)
